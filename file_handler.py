"""
file_handler.py — Universal file type handler.

Supports:
  PDF    → PyMuPDF / pdf2image
  Images → PIL direct (PNG, JPG, BMP, TIFF, WEBP, GIF)
  Word   → python-docx  (DOCX)
  Excel  → openpyxl     (XLSX, XLS via xlrd)
  PPT    → python-pptx  (PPTX)
  Text   → plain read   (TXT, CSV, JSON, XML, HTML, MD)
"""

import os
import io
import logging
from typing import List, Tuple
from PIL import Image

logger = logging.getLogger(__name__)

# ── Optional imports ──────────────────────────────────────────────────────────
try:
    import fitz
    _PYMUPDF = True
except ImportError:
    _PYMUPDF = False

try:
    from pdf2image import convert_from_path
    _PDF2IMAGE = True
except ImportError:
    _PDF2IMAGE = False

try:
    import docx
    _DOCX = True
except ImportError:
    _DOCX = False

try:
    import openpyxl
    _OPENPYXL = True
except ImportError:
    _OPENPYXL = False

try:
    import xlrd
    _XLRD = True
except ImportError:
    _XLRD = False

try:
    from pptx import Presentation
    _PPTX = True
except ImportError:
    _PPTX = False

try:
    import pytesseract
    _TESSERACT = True
except ImportError:
    _TESSERACT = False

# ── File type groups ──────────────────────────────────────────────────────────
PDF_TYPES   = {".pdf"}
IMAGE_TYPES = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp", ".gif"}
WORD_TYPES  = {".docx", ".doc"}
EXCEL_TYPES = {".xlsx", ".xls", ".csv"}
PPT_TYPES   = {".pptx", ".ppt"}
TEXT_TYPES  = {".txt", ".md", ".json", ".xml", ".html", ".htm", ".log", ".yaml", ".yml", ".ini", ".cfg"}

ALL_SUPPORTED = (
    PDF_TYPES | IMAGE_TYPES | WORD_TYPES |
    EXCEL_TYPES | PPT_TYPES | TEXT_TYPES
)

# Streamlit file_uploader type list (without dots)
STREAMLIT_TYPES = [t.lstrip(".") for t in sorted(ALL_SUPPORTED)]


def get_file_category(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext in PDF_TYPES:    return "pdf"
    if ext in IMAGE_TYPES:  return "image"
    if ext in WORD_TYPES:   return "word"
    if ext in EXCEL_TYPES:  return "excel"
    if ext in PPT_TYPES:    return "ppt"
    if ext in TEXT_TYPES:   return "text"
    return "unknown"


# ── Universal: file → images ──────────────────────────────────────────────────

def file_to_images(path: str, dpi: int = 150) -> List[Image.Image]:
    """Convert any supported file into a list of PIL Images (one per page/sheet)."""
    cat = get_file_category(path)

    if cat == "pdf":
        return _pdf_to_images(path, dpi)
    if cat == "image":
        return _image_to_images(path)
    if cat == "word":
        return _word_to_images(path, dpi)
    if cat == "excel":
        return _excel_to_images(path)
    if cat == "ppt":
        return _ppt_to_images(path, dpi)
    if cat == "text":
        return _text_to_images(path)

    logger.warning("Unsupported file type: %s", path)
    return [Image.new("RGB", (800, 1100), (255, 255, 255))]


# ── Universal: file → text ────────────────────────────────────────────────────

def file_to_text(path: str, images: List[Image.Image] = None) -> str:
    """Extract text from any supported file."""
    cat = get_file_category(path)

    if cat == "pdf":
        return _pdf_to_text(path, images or [])
    if cat == "image":
        return _image_to_text(images or file_to_images(path))
    if cat == "word":
        return _word_to_text(path)
    if cat == "excel":
        return _excel_to_text(path)
    if cat == "ppt":
        return _ppt_to_text(path)
    if cat == "text":
        return _plain_text(path)

    return ""


# ── PDF ───────────────────────────────────────────────────────────────────────

def _pdf_to_images(path, dpi):
    if _PYMUPDF:
        try:
            doc = fitz.open(path)
            pages = []
            for page in doc:
                mat = fitz.Matrix(dpi / 72, dpi / 72)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                pages.append(Image.frombytes("RGB", [pix.width, pix.height], pix.samples))
            doc.close()
            return pages
        except Exception as e:
            logger.warning("PyMuPDF render failed: %s", e)
    if _PDF2IMAGE:
        try:
            return convert_from_path(path, dpi=dpi)
        except Exception as e:
            logger.warning("pdf2image failed: %s", e)
    return [Image.new("RGB", (800, 1100), (255, 255, 255))]


def _pdf_to_text(path, images):
    if _PYMUPDF:
        try:
            doc = fitz.open(path)
            text = "\n".join(p.get_text("text") for p in doc).strip()
            doc.close()
            if len(text) > 50:
                return text
        except Exception as e:
            logger.warning("PyMuPDF text failed: %s", e)
    if _TESSERACT and images:
        try:
            return "\n".join(pytesseract.image_to_string(img, config="--psm 6") for img in images).strip()
        except Exception as e:
            logger.warning("Tesseract failed: %s", e)
    return ""


# ── Image ─────────────────────────────────────────────────────────────────────

def _image_to_images(path):
    try:
        img = Image.open(path).convert("RGB")
        # Handle multi-frame (GIF / TIFF)
        frames = []
        try:
            for i in range(getattr(img, "n_frames", 1)):
                img.seek(i)
                frames.append(img.copy().convert("RGB"))
        except EOFError:
            pass
        return frames if frames else [img]
    except Exception as e:
        logger.warning("Image open failed: %s", e)
        return [Image.new("RGB", (800, 1100), (255, 255, 255))]


def _image_to_text(images):
    if _TESSERACT:
        try:
            return "\n".join(pytesseract.image_to_string(img, config="--psm 6") for img in images).strip()
        except Exception as e:
            logger.warning("Tesseract on image failed: %s", e)
    return "[Image file — OCR not available]"


# ── Word (DOCX) ───────────────────────────────────────────────────────────────

def _word_to_images(path, dpi):
    """
    Best-effort: extract text and render it onto a white canvas.
    Full fidelity requires LibreOffice; this gives a readable text render.
    """
    text = _word_to_text(path)
    return _text_string_to_images(text, title="Word Document")


def _word_to_text(path):
    if _DOCX and path.endswith(".docx"):
        try:
            doc = docx.Document(path)
            parts = []
            for para in doc.paragraphs:
                parts.append(para.text)
            # Also grab table text
            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
            return "\n".join(parts).strip()
        except Exception as e:
            logger.warning("python-docx failed: %s", e)
    return "[DOCX extraction unavailable — install python-docx]"


# ── Excel (XLSX / XLS / CSV) ──────────────────────────────────────────────────

def _excel_to_images(path):
    text = _excel_to_text(path)
    return _text_string_to_images(text, title="Spreadsheet")


def _excel_to_text(path):
    ext = os.path.splitext(path)[1].lower()

    if ext == ".csv":
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            logger.warning("CSV read failed: %s", e)
        return ""

    if ext == ".xlsx" and _OPENPYXL:
        try:
            wb = openpyxl.load_workbook(path, data_only=True)
            lines = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                lines.append(f"=== Sheet: {sheet} ===")
                for row in ws.iter_rows(values_only=True):
                    lines.append(" | ".join(str(c) if c is not None else "" for c in row))
            return "\n".join(lines)
        except Exception as e:
            logger.warning("openpyxl failed: %s", e)

    if ext == ".xls" and _XLRD:
        try:
            import xlrd
            wb = xlrd.open_workbook(path)
            lines = []
            for sheet in wb.sheets():
                lines.append(f"=== Sheet: {sheet.name} ===")
                for r in range(sheet.nrows):
                    lines.append(" | ".join(str(sheet.cell_value(r, c)) for c in range(sheet.ncols)))
            return "\n".join(lines)
        except Exception as e:
            logger.warning("xlrd failed: %s", e)

    return "[Excel extraction unavailable — install openpyxl]"


# ── PowerPoint (PPTX) ─────────────────────────────────────────────────────────

def _ppt_to_images(path, dpi):
    text = _ppt_to_text(path)
    return _text_string_to_images(text, title="Presentation")


def _ppt_to_text(path):
    if _PPTX and path.endswith(".pptx"):
        try:
            prs = Presentation(path)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"=== Slide {i} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
            return "\n".join(lines)
        except Exception as e:
            logger.warning("python-pptx failed: %s", e)
    return "[PPTX extraction unavailable — install python-pptx]"


# ── Plain text / CSV / JSON / XML / HTML / MD ─────────────────────────────────

def _plain_text(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        logger.warning("Plain text read failed: %s", e)
        return ""


def _text_to_images(path):
    text = _plain_text(path)
    return _text_string_to_images(text, title=os.path.basename(path))


# ── Helper: render text onto image canvas ─────────────────────────────────────

def _text_string_to_images(text: str, title: str = "", max_lines_per_page: int = 60) -> List[Image.Image]:
    """
    Render a long text string onto one or more A4-ish white image pages.
    Uses PIL ImageDraw — no external dependencies.
    """
    from PIL import ImageDraw, ImageFont

    W, H   = 850, 1100
    MARGIN = 40
    LINE_H = 16
    FONT_SIZE = 13

    try:
        font = ImageFont.truetype("arial.ttf", FONT_SIZE)
    except Exception:
        font = ImageFont.load_default()

    lines = []
    if title:
        lines.append(f"[ {title} ]")
        lines.append("")
    for raw in text.splitlines():
        # Wrap long lines
        while len(raw) > 95:
            lines.append(raw[:95])
            raw = raw[95:]
        lines.append(raw)

    pages = []
    for chunk_start in range(0, max(1, len(lines)), max_lines_per_page):
        chunk = lines[chunk_start: chunk_start + max_lines_per_page]
        img  = Image.new("RGB", (W, H), (255, 255, 255))
        draw = ImageDraw.Draw(img)
        # Light header bar
        draw.rectangle([0, 0, W, 32], fill=(240, 244, 248))
        draw.text((MARGIN, 8), title or "Document", fill=(30, 41, 59), font=font)
        y = 48
        for line in chunk:
            draw.text((MARGIN, y), line, fill=(30, 41, 59), font=font)
            y += LINE_H
        pages.append(img)

    return pages if pages else [Image.new("RGB", (850, 1100), (255, 255, 255))]
